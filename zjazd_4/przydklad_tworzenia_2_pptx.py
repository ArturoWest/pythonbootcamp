import csv
from pptx import Presentation
prs = Presentation()

slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes

title_shape = shapes.title

body_shape = shapes.placeholders[1]

title_shape.text = 'Jakos text'

tf = body_shape.text_frame
tf.text = "zawartosc tekst frame"

with open("report.csv") as f:
    data = cssv.reader(f, delimiter=",")
    for row in data:

    p = tf.add_paragraph()
    p.text = "kobiety"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Przezylo: 70%"
    p.level = 1

prs.save("raport.pptx")